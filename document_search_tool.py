import os
import re
from tkinter import Tk, Frame, Text, Entry, Button, Scrollbar, RIGHT, Y, LEFT, BOTH, TOP, END, ttk, filedialog

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from PIL import Image, ImageTk, ImageGrab
except ImportError:
    Image = ImageTk = ImageGrab = None


class DocumentSearchTool:
    def __init__(self, root_dir="."):
        self.root_dir = os.path.abspath(root_dir)
        self.matches = []
        self.root = Tk()
        self.root.title("Document Search Tool")
        self._setup_ui()

    def _setup_ui(self):
        left_frame = Frame(self.root)
        left_frame.pack(side=LEFT, fill=Y)

        right_frame = Frame(self.root)
        right_frame.pack(side=LEFT, fill=BOTH, expand=True)

        self.tree = ttk.Treeview(left_frame)
        self.tree.pack(side=LEFT, fill=Y, expand=True)
        scroll = Scrollbar(left_frame, command=self.tree.yview)
        scroll.pack(side=RIGHT, fill=Y)
        self.tree.configure(yscrollcommand=scroll.set)
        self._populate_tree(self.root_dir)
        self.tree.bind("<<TreeviewSelect>>", self._on_tree_select)

        search_frame = Frame(right_frame)
        search_frame.pack(side=TOP, fill="x")

        self.search_entry = Entry(search_frame)
        self.search_entry.pack(side=LEFT, fill=BOTH, expand=True)
        Button(search_frame, text="Search", command=self.search).pack(side=LEFT)

        self.result_box = Text(right_frame)
        self.result_box.pack(fill=BOTH, expand=True)

    def _populate_tree(self, path, parent=""):
        for name in os.listdir(path):
            full = os.path.join(path, name)
            node = self.tree.insert(parent, END, text=name, open=False)
            if os.path.isdir(full):
                self._populate_tree(full, node)

    def _on_tree_select(self, event):
        sel = self.tree.selection()
        if sel:
            node = sel[0]
            path = self._get_full_path(node)
            if os.path.isdir(path):
                self.root_dir = path

    def _get_full_path(self, node):
        parts = []
        while node:
            parts.append(self.tree.item(node, "text"))
            node = self.tree.parent(node)
        return os.path.join(*reversed(parts))

    def search(self):
        query = self.search_entry.get().strip()
        if not query:
            return
        self.matches.clear()
        for root, _, files in os.walk(self.root_dir):
            for f in files:
                full = os.path.join(root, f)
                text = self._extract_text(full)
                if text and query.lower() in text.lower():
                    self.matches.append((full, text))
        self._display_results(query)

    def _display_results(self, query):
        self.result_box.delete(1.0, END)
        for idx, (path, text) in enumerate(self.matches, 1):
            self.result_box.insert(END, f"{idx}. {path}\n")
            start = text.lower().find(query.lower())
            snippet = text[start:start+200] if start != -1 else ""
            self.result_box.insert(END, f"    ...{snippet}...\n\n")
        self.result_box.tag_remove('highlight', '1.0', END)
        if query:
            start = '1.0'
            while True:
                pos = self.result_box.search(query, start, END, nocase=True)
                if not pos:
                    break
                end_pos = f"{pos}+{len(query)}c"
                self.result_box.tag_add('highlight', pos, end_pos)
                start = end_pos
            self.result_box.tag_config('highlight', background='yellow')

    def _extract_text(self, path):
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext == '.pdf' and fitz:
                doc = fitz.open(path)
                text = "\n".join(page.get_text() for page in doc)
                doc.close()
                return text
            elif ext == '.docx' and Document:
                doc = Document(path)
                return "\n".join(p.text for p in doc.paragraphs)
            elif ext == '.pptx' and Presentation:
                prs = Presentation(path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            elif ext in ('.xlsx', '.xls') and load_workbook:
                wb = load_workbook(path)
                text = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        text.append("\t".join([str(c) if c is not None else '' for c in row]))
                return "\n".join(text)
            elif ext == '.txt':
                with open(path, 'r', errors='ignore') as f:
                    return f.read()
        except Exception:
            return ''
        return ''

    def run(self):
        self.root.mainloop()


if __name__ == '__main__':
    root = filedialog.askdirectory(title="Select Root Directory")
    if root:
        app = DocumentSearchTool(root)
        app.run()
